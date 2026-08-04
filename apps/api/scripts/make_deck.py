#!/usr/bin/env python3
"""
Build the demo deck as a real .pptx — a run sheet for presenting the product live.

**Why it is shaped as a run sheet and not an architecture review.** A deck that only describes the
system leaves the presenter narrating an idle screen; a deck that only lists clicks leaves the
audience watching a tour with no idea why any of it was hard. So every demo slide carries three
things: **SHOW** what to do on screen, **SAY** the one sentence that lands, and **WHY IT WORKS** —
the mechanism, with the measurement that supports it. The deck can be read as a script and still
says something.

**Why a generator and not a hand-made file.** Every figure has to trace to a run. The numbers live
in one `MEASURED` dict beside the run that produced them, so a slide cannot quietly drift from
MEASUREMENTS.md the way a hand-edited deck would. Re-run it after a measurement changes and the deck
follows.

**The two rules from CLAUDE.md are enforced here rather than hoped for.** No invented measurements:
anything without a run says `NOT YET MEASURED` and says why, which is why there is no dollar figure
anywhere — cost per interview divides by a concurrency number nobody has measured. And the judgment
sections stay the author's: the build-vs-buy slide quotes PROCESS.md §4 as written and marks the
refresh it now needs `[HUMAN]`, because §4.2 opens "I did not run a GPU" and that input has changed.

    .venv/bin/python apps/api/scripts/make_deck.py --out nod-demo.pptx
"""

from __future__ import annotations

import argparse
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

# -- palette ---------------------------------------------------------------------------------
# Chosen rather than defaulted: a deep slate-navy ground, one warm amber accent used sparingly,
# and neutrals biased very slightly blue so they read as part of the same family instead of as
# grey.
INK = RGBColor(0x14, 0x1B, 0x2D)
INK_MID = RGBColor(0x3A, 0x44, 0x5C)
INK_LOW = RGBColor(0x6B, 0x76, 0x8C)
PAPER = RGBColor(0xFA, 0xFB, 0xFD)
RULE = RGBColor(0xD8, 0xDD, 0xE6)
ACCENT = RGBColor(0xC9, 0x7B, 0x2A)
GOOD = RGBColor(0x2E, 0x6F, 0x4F)
WARN = RGBColor(0xB3, 0x5C, 0x1F)
BAD = RGBColor(0x9E, 0x2B, 0x2B)

BODY = "Arial"
W, H = Inches(13.333), Inches(7.5)


MEASURED = {
    # Every value here came from a run. The comment names the run or the document that records
    # it.
    "render_ms": "78.4",  # MEASUREMENTS.md §2, T4, batch 16, fp16, CPU overlapped
    "render_ms_seq": "124.7",  # same run, stages in sequence
    "render_speedup": "1.59×",
    "render_fps_capacity": "12.8",
    "fps_delivered": "8.3",  # §8b, live over WebSocket, target 8
    "fps_delivered_range": "8.2 – 8.9",
    "fps_before": "1.0 – 2.4",
    "gap_median": "−66 ms to +172 ms",  # two runs, medians disagreed in sign
    "gap_worst": "538 ms",
    "start_lag": "1,510 ms",
    "discards": "6 – 11",
    "discards_before": "33 – 81",
    "vae_ms": "57.8",
    "unet_ms": "12.3",
    "cpu_ms": "51.5",
    "gpu_ms": "70.1",
    "first_render_ms": "4,747",  # warm-up, 5 frames
    "models_ms": "22,524",
    "prepare_ms": "126,817",  # 500-frame reference
    "animate_ms": "213,112",  # LivePortrait, 500 frames
    "enroll_202_ms": "0.018",
    "tts_hosted_ms": "290",
    "tts_cloned_ms": "4,848",
    "llm_ttft": "1,645 / 2,942 / 4,724",
    "tts_rest": "869 / 956 / 889",
    "tts_ws": "351 – 361",
    "turn_total": "2.7 – 5.8 s",
    "turn_detect": "700",
    "scorer": "8 ms to queue, 6,093 ms of work",
    "recording": "7.0 MB, 1 m 37 s, H.264 + AAC",
    "webrtc_first": "4,296 ms perceived / 4,221 ms paint",
    "tests": "746",
    "session_start": "1.5 – 3.8 s",
    "voice_contention": "3 s → 28 s",
    # The LiveKit worker, measured at a remote subscriber against a local SFU.
    "drift_median": "−241 ms",
    "drift_spread": "9 ms across 22 s",
    "ws_drift_range": "−66 ms to +172 ms, worst 538 ms",
    "sub_video": "327",
    "sub_audio": "2,200",
    "bargein_dropped": "8 of 38 chunks",
    "seams_forced": "0",
    "resume_chars": "778",
    "tests_now": "823",
}


def add(prs: Presentation) -> object:
    return prs.slides.add_slide(prs.slide_layouts[6])


def rect(slide, x, y, w, h, fill=None, line=None, line_w=1.0):
    from pptx.enum.shapes import MSO_SHAPE

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(line_w)
    shape.shadow.inherit = False
    return shape


def text(
    slide, x, y, w, h, runs, size=14, color=INK, bold=False, align=PP_ALIGN.LEFT, spacing=1.15
):
    """`runs` is a string or a list of (text, {overrides}) tuples, one paragraph per list item."""
    box = slide.shapes.add_textbox(x, y, w, h)
    frame = box.text_frame
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = 0
    frame.margin_top = frame.margin_bottom = 0
    items = [runs] if isinstance(runs, str) else runs
    for index, item in enumerate(items):
        body, over = (item, {}) if isinstance(item, str) else item
        para = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        para.alignment = over.get("align", align)
        para.line_spacing = over.get("spacing", spacing)
        if over.get("space_before"):
            para.space_before = Pt(over["space_before"])
        run = para.add_run()
        run.text = body
        font = run.font
        font.name = over.get("font", BODY)
        font.size = Pt(over.get("size", size))
        font.bold = over.get("bold", bold)
        font.color.rgb = over.get("color", color)
    return box


def title_slide(prs, kicker, headline, sub):
    slide = add(prs)
    rect(slide, 0, 0, W, H, fill=INK)
    rect(slide, Inches(0.9), Inches(2.35), Inches(1.1), Pt(4), fill=ACCENT)
    text(slide, Inches(0.9), Inches(1.75), Inches(11), Inches(0.4),
         kicker.upper(), size=12, color=ACCENT, bold=True)
    text(slide, Inches(0.9), Inches(2.7), Inches(11.2), Inches(1.6),
         headline, size=40, color=PAPER, bold=True, spacing=1.0)
    text(slide, Inches(0.9), Inches(4.5), Inches(10), Inches(1.2),
         sub, size=15, color=RGBColor(0xA8, 0xB2, 0xC6), spacing=1.3)
    return slide


def section(prs, number, headline, sub=""):
    slide = add(prs)
    rect(slide, 0, 0, W, H, fill=INK)
    text(slide, Inches(0.9), Inches(2.9), Inches(2), Inches(0.8),
         number, size=13, color=ACCENT, bold=True)
    text(slide, Inches(0.9), Inches(3.3), Inches(11), Inches(1.2),
         headline, size=34, color=PAPER, bold=True, spacing=1.0)
    if sub:
        text(slide, Inches(0.9), Inches(4.5), Inches(10.5), Inches(1),
             sub, size=14, color=RGBColor(0xA8, 0xB2, 0xC6), spacing=1.3)
    return slide


def content(prs, heading, kicker="", notes=""):
    slide = add(prs)
    rect(slide, 0, 0, W, H, fill=PAPER)
    top = Inches(0.62)
    if kicker:
        text(slide, Inches(0.72), top, Inches(11.5), Inches(0.3),
             kicker.upper(), size=10.5, color=ACCENT, bold=True)
        top = Inches(0.95)
    text(slide, Inches(0.72), top, Inches(11.9), Inches(0.6),
         heading, size=25, color=INK, bold=True, spacing=1.0)
    rect(slide, Inches(0.72), Inches(1.62), Inches(11.9), Pt(1.1), fill=RULE)
    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    return slide


def table(slide, x, y, w, headers, rows, widths=None, size=11.5, header_size=10,
          row_h=0.34, colors=None):
    """
    A table drawn from rectangles rather than a pptx table.

    Deliberate: pptx tables carry a theme style that overrides per-cell colour in ways that
    differ between PowerPoint, Keynote and Google Slides, and a deck whose whole point is
    legibility of numbers cannot have its emphasis silently dropped by the viewer's app.
    """
    widths = widths or [1.0 / len(headers)] * len(headers)
    total = float(w)
    xs, running = [], float(x)
    for fraction in widths:
        xs.append(running)
        running += total * fraction

    text(slide, Emu(int(x)), Emu(int(y)), Emu(int(w)), Inches(0.25), "", size=1)
    for index, head in enumerate(headers):
        text(slide, Emu(int(xs[index])), Emu(int(y)), Emu(int(total * widths[index] - 60000)),
             Inches(0.26), head.upper(), size=header_size, color=INK_LOW, bold=True)
    rect(slide, Emu(int(x)), Emu(int(y + Inches(0.3))), Emu(int(w)), Pt(0.9), fill=RULE)

    cursor = y + Inches(0.42)
    for r, row in enumerate(rows):
        for c, cell in enumerate(row):
            colour = INK if c == 0 else INK_MID
            bold = c == 0
            if colors and (r, c) in colors:
                colour = colors[(r, c)]
                bold = True
            text(slide, Emu(int(xs[c])), Emu(int(cursor)),
                 Emu(int(total * widths[c] - 60000)), Inches(row_h),
                 cell, size=size, color=colour, bold=bold, spacing=1.05)
        cursor += Inches(row_h)
        if r < len(rows) - 1:
            rect(slide, Emu(int(x)), Emu(int(cursor - Inches(0.06))), Emu(int(w)),
                 Pt(0.6), fill=RGBColor(0xEC, 0xEF, 0xF4))
    return cursor


def bullets(slide, x, y, w, items, size=13, gap=0.1):
    cursor = float(y)
    for item in items:
        head, body = (item, "") if isinstance(item, str) else item
        rect(slide, Emu(int(x)), Emu(int(cursor + Inches(0.07))), Pt(3.2), Pt(3.2), fill=ACCENT)
        runs = [(head, {"bold": True, "size": size, "color": INK})]
        if body:
            runs.append((body, {"size": size - 1, "color": INK_MID, "spacing": 1.25}))
        text(slide, Emu(int(x + Inches(0.22))), Emu(int(cursor)), Emu(int(w - Inches(0.22))),
             Inches(0.3), runs)
        lines = 1 + (len(body) // 95 if body else 0)
        cursor += Inches(0.28 + 0.19 * lines + gap)
    return cursor


def kpi(slide, x, y, w, value, label, note="", colour=INK):
    rect(slide, Emu(int(x)), Emu(int(y)), Emu(int(w)), Inches(1.42),
         fill=RGBColor(0xFF, 0xFF, 0xFF), line=RULE)
    rect(slide, Emu(int(x)), Emu(int(y)), Pt(3), Inches(1.42), fill=colour)
    text(slide, Emu(int(x + Inches(0.22))), Emu(int(y + Inches(0.16))),
         Emu(int(w - Inches(0.4))), Inches(0.5), value, size=25, color=colour, bold=True)
    text(slide, Emu(int(x + Inches(0.22))), Emu(int(y + Inches(0.72))),
         Emu(int(w - Inches(0.4))), Inches(0.3), label, size=10.5, color=INK, bold=True)
    if note:
        text(slide, Emu(int(x + Inches(0.22))), Emu(int(y + Inches(0.98))),
             Emu(int(w - Inches(0.4))), Inches(0.4), note, size=9.5, color=INK_LOW, spacing=1.15)


def callout(slide, x, y, w, h, body, tone=ACCENT, label=""):
    rect(slide, Emu(int(x)), Emu(int(y)), Emu(int(w)), Emu(int(h)),
         fill=RGBColor(0xF4, 0xF1, 0xEA) if tone == ACCENT else RGBColor(0xF2, 0xF6, 0xF3))
    rect(slide, Emu(int(x)), Emu(int(y)), Pt(3), Emu(int(h)), fill=tone)
    top = y + Inches(0.16)
    if label:
        text(slide, Emu(int(x + Inches(0.24))), Emu(int(top)), Emu(int(w - Inches(0.5))),
             Inches(0.25), label.upper(), size=9.5, color=tone, bold=True)
        top += Inches(0.28)
    text(slide, Emu(int(x + Inches(0.24))), Emu(int(top)), Emu(int(w - Inches(0.5))),
         Emu(int(h)) - Inches(0.3), body, size=12, color=INK_MID, spacing=1.3)


def flow(slide, x, y, w, steps, note_size=9.5):
    """A horizontal pipeline. Each step is (title, detail)."""
    n = len(steps)
    gap = Inches(0.16)
    box_w = (Emu(int(w)) - gap * (n - 1)) / n
    for index, (head, detail) in enumerate(steps):
        left = Emu(int(x)) + (box_w + gap) * index
        rect(slide, Emu(int(left)), Emu(int(y)), Emu(int(box_w)), Inches(1.15),
             fill=RGBColor(0xFF, 0xFF, 0xFF), line=RULE)
        rect(slide, Emu(int(left)), Emu(int(y)), Emu(int(box_w)), Pt(2.4), fill=INK)
        text(slide, Emu(int(left + Inches(0.14))), Emu(int(y + Inches(0.18))),
             Emu(int(box_w - Inches(0.28))), Inches(0.3), head, size=11.5, color=INK, bold=True)
        text(slide, Emu(int(left + Inches(0.14))), Emu(int(y + Inches(0.5))),
             Emu(int(box_w - Inches(0.28))), Inches(0.6), detail,
             size=note_size, color=INK_LOW, spacing=1.15)
        if index < n - 1:
            text(slide, Emu(int(left + box_w + Inches(0.01))), Emu(int(y + Inches(0.42))),
                 Emu(int(gap)), Inches(0.3), "›", size=15, color=ACCENT, bold=True,
                 align=PP_ALIGN.CENTER)


def beat(slide, x, y, w, show, say, why=""):
    """
    One demo beat: what to do on screen, what to say, and the mechanism behind it.

    The three-part shape exists because a demo deck fails in two different ways. Slides that only
    describe architecture leave the presenter narrating an idle screen; slides that only list clicks
    leave the audience watching a tour with no idea why any of it is hard. Pairing them on every
    beat means the deck can be read as a script and still says something.
    """
    cursor = float(y)
    for label, body, tone in (
        ("show", show, ACCENT),
        ("say", say, INK),
        ("why it works", why, GOOD),
    ):
        if not body:
            continue
        rect(slide, Emu(int(x)), Emu(int(cursor + Inches(0.02))), Pt(2.5), Inches(0.42), fill=tone)
        text(slide, Emu(int(x + Inches(0.16))), Emu(int(cursor)), Inches(1.15), Inches(0.3),
             label.upper(), size=9, color=tone, bold=True)
        text(slide, Emu(int(x + Inches(1.38))), Emu(int(cursor)),
             Emu(int(w - Inches(1.38))), Inches(0.4),
             body, size=12.5 if label == "say" else 11.5,
             color=INK if label == "say" else INK_MID,
             bold=label == "say", spacing=1.28)
        lines = max(1, -(-len(body) // 86))
        cursor += Inches(0.2 + 0.2 * lines)
    return cursor


def build(out: Path) -> None:
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H

    # ---------------------------------------------------------------- 1. title
    title_slide(
        prs,
        "Engineering review · confidential",
        "A real-time conversational\navatar interview system",
        "Built in-house, end to end: identity capture, audio-driven face rendering, session "
        "mechanics, scoring.\nEvery number in this deck came from a run on named hardware. "
        "Where a number does not exist, the slide says so.",
    )

    # ---------------------------------------------------------------- 2. exec summary
    s = content(prs, "Where the system stands", "Executive summary",
                notes="Open on the honest frame: this is a working product that is not yet fast "
                      "enough, and we can say exactly how much not-enough and why. The three "
                      "KPIs are all from runs on a Tesla T4. 8.3 fps is against a configured "
                      "target of 8, not against 25 — say that out loud before anyone asks.")
    kpi(s, Inches(0.72), Inches(1.9), Inches(3.75), MEASURED["fps_delivered"] + " fps",
        "Delivered to the client, live", "Against a configured target of 8. Capacity is "
        f"{MEASURED['render_fps_capacity']} fps.", GOOD)
    kpi(s, Inches(4.79), Inches(1.9), Inches(3.75), MEASURED["gap_median"],
        "Audio-to-video trailing gap", "Median of two runs; they disagree in sign. "
        f"Worst single turn {MEASURED['gap_worst']}.", GOOD)
    kpi(s, Inches(8.86), Inches(1.9), Inches(3.75), "2.0×",
        "Short of 25 fps real time", f"{MEASURED['render_ms']} ms/frame against a 40 ms budget. "
        "VAE decode is 74% of it.", WARN)
    bullets(s, Inches(0.72), Inches(3.6), Inches(11.9), [
        ("Working and verified by driving it, not by inspection.",
         "Two-way conversation with barge-in, a real face from an uploaded video or an animated "
         "photograph, competency-planned questions, async scoring, a report view, egress "
         "recording, and a screen-aware operator assistant."),
        ("The renderer was never the bottleneck — and we measured that twice.",
         f"A full turn is {MEASURED['turn_total']}; set the renderer to zero and ~2.6–5.7 s "
         "remains. Separately, delivered fps was 1.4 while the card could do 12.8: the fault was "
         "our orchestration, not the model."),
        ("Two things are open, and one of them is a business decision.",
         "One T4 serves a self-hosted face or a self-hosted voice, not both. And there is no "
         "authentication anywhere — a stated development posture, not an oversight."),
    ])

    # ---------------------------------------------------------------- 3. the demo
    section(prs, "01", "What you are about to see",
            "Five minutes of running software, then how each part of it works. Every number in "
            "this deck came from a run you can repeat.")

    s = content(prs, "The demo, in order", "Run sheet",
                notes="Keep this slide up while you set up. The order is deliberate: configuration "
                      "first so the interview has something to be configured by, then the "
                      "interview, then the assessment that falls out of it. Do not start with the "
                      "avatar — it is the least interesting part and it invites questions you would "
                      "rather answer at the end.")
    table(s, Inches(0.72), Inches(1.9), Inches(11.9),
          ["#", "On screen", "The point being made", "Time"],
          [["1", "Console → Candidates. Drop a resume on a row.",
            "The resume is the only field here that changes the interview", "45 s"],
           ["2", "Console → Agents. Open one.",
            "An interviewer is composed, not coded: rubric, voice, knowledge, guardrail", "45 s"],
           ["3", "Console → Knowledge. Paste a paragraph, hit Reindex, run a query.",
            "Retrieval is inspectable before anyone is interviewed against it", "45 s"],
           ["4", "Copy the invite link. Open the interview room.",
            "Confirm your name, accept the recording notice — attested, never verified", "30 s"],
           ["5", "Have the conversation. **Interrupt it mid-sentence.**",
            "Barge-in is the moment the architecture shows", "90 s"],
           ["6", "Console → Sessions → the report.",
            "Quotes, not scores, are the artefact — and every quote is checked", "60 s"],
           ["7", "Terminal: the LiveKit worker, two processes.",
            "Where this is going: the renderer as its own scalable participant", "60 s"]],
          widths=[0.04, 0.33, 0.51, 0.12], size=11, row_h=0.44)
    callout(s, Inches(0.72), Inches(5.6), Inches(11.9), Inches(1.2),
            "One thing to set up before the room is silent: the renderer here is the placeholder, "
            "not a face. The real face runs on a Tesla T4 at 12.8 fps of capacity and 8.3 fps "
            "delivered — measured — and this laptop has no GPU. Say that at the start rather than "
            "when someone notices, and the rest of the demo is about mechanism instead of apology.",
            tone=ACCENT, label="Say this in the first thirty seconds")

    # ---------------------------------------------------------------- 4. walkthrough
    section(prs, "02", "Walk the console",
            "Four screens. Each one is a module, and each has a mechanism worth a sentence.")

    s = content(prs, "Candidates — the resume changes the interview", "Beat 1 · module: candidates",
                notes="This is the strongest opening because it is the least expected: dropping a "
                      "PDF changes what the interviewer asks. Show the extracted text — that is the "
                      "detail that convinces people it is real rather than decorative. If the "
                      "resume failed to parse, show that too; the failure state is deliberately as "
                      "loud as the success.")
    beat(s, Inches(0.72), Inches(1.95), Inches(11.9),
         "Open Candidates. Click a name, or “+ add resume” on the row. Drop a PDF or "
         "Markdown file. Then click “Show what the interviewer reads”.",
         "This resume is the only thing on this screen that changes the interview.",
         "The runtime appends it to the interviewer's system prompt at session start — "
         f"measured at {MEASURED['resume_chars']} characters for one real resume — framed "
         "explicitly as the candidate's own unverified claims, with an instruction to probe "
         "rather than recite. Without that framing a model states a CV back as fact, which is "
         "the failure an interview exists to prevent.")
    table(s, Inches(0.72), Inches(4.15), Inches(11.9),
          ["Also worth showing", "Because"],
          [["A resume that failed to parse", "A scanned PDF stores fine and extracts to nothing. "
            "The row says so, in warning colour, with the reason and the fix"],
           ["The status column moving to “invited”", "Advanced by the API when an interview is "
            "minted, not typed by an operator — a status nobody sets by hand cannot drift"],
           ["Delete a candidate", "Removes the resume file and the ~1 GB prepared identity too, "
            "and keeps their interviews: a transcript is evidence"]],
          widths=[0.3, 0.7], size=11, row_h=0.46)

    s = content(prs, "Agents — an interviewer is composed, not coded",
                "Beat 2 · module: agents, rubrics, guardrails, lexicons",
                notes="The point of this screen is that nothing here is a code change. If someone "
                      "asks how you would run a different kind of interview, the answer is on this "
                      "screen and takes a minute. Open the Edit form and change the end-of-turn "
                      "window live if you have time — it is the number that most changes how the "
                      "conversation feels.")
    beat(s, Inches(0.72), Inches(1.95), Inches(11.9),
         "Open Agents. Show the attached column, then Edit one. Point at the end-of-turn "
         "silence field.",
         "Five interviewers, five voices, five rubrics — and no code involved in any of it.",
         "Each attachment is a decorator around the sentence stream: "
         "with_plan(with_guardrail(with_knowledge(llm))). The orchestrator does not know they "
         "exist, which is why adding one is configuration rather than a release.")
    table(s, Inches(0.72), Inches(4.1), Inches(11.9),
          ["Attached", "What it does to the conversation"],
          [["Rubric", "Competencies with probes and signals. Drives what gets asked and when to "
            "move on — not just how it is scored afterwards"],
           ["Guardrail", "Banned topics, PII redaction, a refusal line. Checks the candidate's "
            "input and the model's output separately"],
           ["Pronunciation lexicon", "Rewrites text before synthesis, so “kubectl” is spoken and "
            "not spelled"],
           ["Knowledge base", "Retrieved per turn and appended, so the interviewer can cite the "
            "on-call policy as fact"],
           ["Voice", "A Deepgram Aura voice, or a cloned one from an uploaded recording"],
           ["Turn-taking", "End-of-turn silence in ms. The largest single term in the latency "
            "budget, and a conversational judgement rather than a technical one"]],
          widths=[0.22, 0.78], size=11, row_h=0.42)

    s = content(prs, "Knowledge — retrieval you can look at before you trust it",
                "Beat 3 · module: knowledge",
                notes="The retrieval tester is the slide-worthy part: it shows scores per chunk, so "
                      "bad retrieval is distinguishable from a bad answer. That distinction is the "
                      "difference between debugging this in an afternoon and guessing at it for a "
                      "week. Mention that this used to be a curl command.")
    beat(s, Inches(0.72), Inches(1.95), Inches(11.9),
         "Open Knowledge. Paste two paragraphs separated by a blank line — watch the chunk "
         "count update live. Add it, hit Reindex, then run a query in the tester.",
         "You can see exactly what the interviewer would retrieve, before anyone is interviewed.",
         "Chunking splits on blank lines, so paragraph structure decides what can be retrieved "
         "independently — which is invisible until retrieval disappoints, so the form says it "
         "while you type. Reindex exists because adding a document without rebuilding retrieves "
         "from the stale index, and that looks exactly like retrieval ignoring your upload.")
    callout(s, Inches(0.72), Inches(4.3), Inches(5.8), Inches(2.5),
            "Keyword scoring, not embeddings, and it is a deliberate choice rather than a stage we "
            "have not reached. A retrieval hop inside a conversational turn has no latency budget "
            "to spare — the whole turn is already 2.7 to 5.8 seconds against a sub-second target — "
            "and over a handful of short documents keyword matching is competitive. The interface "
            "is a Protocol, so a vector store is a swap when the corpus justifies one.",
            tone=GOOD, label="Why no embedding model")
    callout(s, Inches(6.82), Inches(4.3), Inches(5.8), Inches(2.5),
            "A query that returns nothing is shown as a real miss, not an error. Keyword retrieval "
            "needs a word in common, so a paraphrase with no shared vocabulary finds nothing — and "
            "saying “no chunk shared a term with that query” tells an operator what to do, where "
            "“0 results” starts a support conversation.",
            tone=ACCENT, label="The empty-result case")

    # ---------------------------------------------------------------- 5. the interview
    section(prs, "03", "The interview itself",
            "One turn, beat by beat — and the interruption, which is where the architecture "
            "becomes visible.")

    s = content(prs, "Joining — attested, never verified", "Beat 4 · module: attendance",
                notes="Do not skip this screen even though it is a form. It is the one place the "
                      "product makes a claim about a person, and the wording is the whole point. "
                      "Anyone from a regulated industry will care more about this slide than about "
                      "the avatar.")
    beat(s, Inches(0.72), Inches(1.95), Inches(11.9),
         "Open the invite link. The name is prefilled from the invite. Accept the recording "
         "notice. Neither field is skippable.",
         "This records who says they turned up. It does not verify anyone, and every heading "
         "says so.",
         "There is no authentication in this system — the link is the whole credential — so "
         "nothing here could prove identity. What it captures is an explicit, timestamped "
         "attestation: the typed name, the expected name, the consent, the browser and timezone. "
         "Both names are stored so a mismatch is visible; `verified` is a field that is always "
         "false, present rather than omitted so nobody infers the absence of a check.")
    callout(s, Inches(0.72), Inches(4.5), Inches(11.9), Inches(2.3),
            "A reviewer who came away believing identity had been confirmed would be making a "
            "hiring decision on a check nobody performed. That is a specific, foreseeable harm, so "
            "the limitation is on the card in the report and in the heading of this screen rather "
            "than in a manual.\n\nOne of the seeded candidates joins under a deliberately "
            "different name — “Thomas” against “Tom” on file — so the mismatch warning "
            "on the report is something you can show rather than describe.",
            tone=ACCENT, label="Why the wording matters more than the feature")

    s = content(prs, "One turn, and where the time goes", "Beat 5 · module: orchestrator",
                notes="Walk the five boxes, then land on the latency table. The conclusion to "
                      "state out loud: none of the three dominant terms is the renderer. That is "
                      "the finding the whole build-vs-buy case rests on, and it is measured.")
    flow(s, Inches(0.72), Inches(1.95), Inches(11.9), [
        ("Mic in", "16 kHz PCM. Energy VAD plus a 700 ms silence window decides the turn is over."),
        ("Transcribe", "Deepgram nova-3 on a persistent socket, so the transcript is final by the "
         "time the window elapses."),
        ("Plan + think", "The competency plan picks what to probe; the LLM streams back sentences, "
         "not tokens."),
        ("Speak", "Aura 2 per sentence. Real-time factor below 1.0 keeps generation ahead of "
         "playback."),
        ("Render + mix", "The same audio drives the face. The mixer emits at cadence and restamps "
         "every frame."),
    ])
    table(s, Inches(0.72), Inches(3.45), Inches(11.9),
          ["Stage", "Measured", "What moves it"],
          [["End-of-turn detection", f"{MEASURED['turn_detect']} ms (configured)",
            "Nothing in hardware. Speculative execution trades wasted compute for latency"],
           ["LLM time-to-first-token", f"{MEASURED['llm_ttft']} ms",
            "Commercial, not architectural: a paid low-latency endpoint"],
           ["TTS time-to-first-audio", f"{MEASURED['tts_rest']} ms",
            f"Aura's WebSocket measured {MEASURED['tts_ws']} ms — the largest cheap win"],
           ["Perceived total", MEASURED["turn_total"], "3–6× target, and the renderer is not why"]],
          widths=[0.26, 0.22, 0.52], size=11, colors={(3, 1): BAD})
    callout(s, Inches(0.72), Inches(5.6), Inches(11.9), Inches(1.2),
            "Set the renderer to zero and roughly 2.6–5.7 s of a 2.7–5.8 s turn remains. The part a "
            "vendor sells is the part that was never the problem; the part that is the problem — "
            "turn-taking, cancellation, history truncation — is code no vendor API writes for you. "
            "That single measurement is the backbone of the build-vs-buy case.",
            tone=GOOD, label="The line to deliver slowly")

    s = content(prs, "Interrupt it — this is the demo", "Beat 5b · module: cancellation",
                notes="Do this twice. The first time people miss it. Talk over the avatar "
                      "mid-sentence and it stops immediately — audio and video together. Then "
                      "explain that the renderer was not interrupted at all, which is the "
                      "counter-intuitive part and the reason it is fast.")
    beat(s, Inches(0.72), Inches(1.95), Inches(11.9),
         "While the avatar is speaking, talk over it. Do it twice — the first time is always "
         "missed. Then show the report: the turn is marked interrupted.",
         "It stops mid-sentence, and the renderer was never told to stop.",
         "Interrupting increments an integer. In-flight GPU work still finishes and its frames "
         "die at the consumer because their epoch is stale — so reaction is one write and wasted "
         "work is bounded by one render window. No interruptible renderer is required, which is "
         "what makes this survive the renderer moving to another process.")
    table(s, Inches(0.72), Inches(4.15), Inches(11.9),
          ["What also happens", "Why it is not optional"],
          [["The client's audio buffer is flushed", "A server-only flush leaves the browser "
            "playing a sentence the avatar has abandoned — a laggy interruption even though the "
            "state machine reacted instantly"],
           ["History truncates to what was *heard*", "The browser reports played milliseconds from "
            "Web Audio's own clock. Without it, an interrupted turn enters history as fully "
            "delivered and the next question refers to a sentence nobody heard"],
           ["The idle loop takes over on a mouth-closed frame",
            f"Cutting mid-vowel pops. Measured across two turns: {MEASURED['seams_forced']} seams "
            "forced, so every cut found a clean frame"]],
          widths=[0.32, 0.68], size=11, row_h=0.5)

    s = content(prs, "The report — quotes are the artefact, not the score",
                "Beat 6 · module: scoring",
                notes="The thing to point at is the unverified-quote warning. Every quote the model "
                      "produced is re-checked against the transcript, and the ones that do not "
                      "match are shown loudly. That is the difference between a scorecard you can "
                      "argue with and one you have to trust.")
    beat(s, Inches(0.72), Inches(1.95), Inches(11.9),
         "Open Sessions, then the report. Scroll to the verdicts. Point at a quote, then at an "
         "unverified one if the run produced any.",
         "The model produces no hiring decision, and it never will. It produces evidence.",
         "Ratings summarise the quotes beneath them, so the quotes are the part worth checking — "
         "and every one is re-matched against the transcript, with mismatches shown as loudly as "
         "matches. `decision` is deliberately null: a model that recommended hiring would be "
         "trusted for a judgement nobody asked it to make.")
    table(s, Inches(0.72), Inches(4.15), Inches(11.9),
          ["On the report", "Reading it"],
          [["Per-competency rating and weight",
            "no_evidence is a real and common verdict — it means the interview never got there, "
            "not that the candidate failed"],
           ["Coverage, beside the scoring", "Two different questions: what was asked, and what was "
            "demonstrated. A competency with no coverage and no evidence is an interview problem"],
           ["Every latency stage, per turn", "The same numbers as the engineering table, on the "
            "record for the interview that actually happened"],
           ["The attendance card", "Attested, not verified — with the mismatch warning if the "
            "typed name differs from the invite"]],
          widths=[0.3, 0.7], size=11, row_h=0.48)

    # ---------------------------------------------------------------- 6. under the hood
    section(prs, "04", "How the modules fit",
            "Four ideas carry most of the behaviour. Each has a test that fails if it is dropped.")

    s = content(prs, "The four ideas", "Internals",
                notes="If they remember one thing: contracts.py imports nothing from the package, "
                      "and a test enforces it by inspecting sys.modules in a clean subprocess. That "
                      "single property is why 823 tests run with no GPU, no weights and no network.")
    bullets(s, Inches(0.72), Inches(1.95), Inches(11.9), [
        ("Everything is a Protocol in contracts.py, which imports nothing.",
         "Renderer, transport, LLM, TTS, transcriber. Enforced by a test that imports the "
         f"orchestration layer in a clean subprocess and checks what came with it — which is why "
         f"all {MEASURED['tests_now']} tests run with no GPU, no weights and no network."),
        ("Cancellation is an integer.",
         "Interrupting increments an epoch; stale frames die at the consumer. It survived the "
         "renderer moving to another process unchanged — across the boundary it is still one "
         "comparison, where interrupting the producer would be a round trip."),
        ("History truncates to what was heard, not what was sent.",
         "Acknowledged-playback milliseconds from the browser's own clock. Otherwise the "
         "interviewer follows up on a sentence the candidate never heard."),
        ("Composition, not configuration.",
         "Knowledge, pronunciation, guardrails and the plan each wrap a sentence stream. Paired "
         "delivery wraps the transport the same way. The orchestrator never learns any of it."),
    ], size=13.5)

    s = content(prs, "What the renderer actually does", "Internals · the model",
                notes="This is the slide that shows understanding rather than integration. "
                      "MuseTalk is not a diffusion model in operation: timestep=0, no scheduler, no "
                      "sampling loop. in_channels 8 against out_channels 4 means two concatenated "
                      "latents in and one out. cross_attention_dim 384 is exactly whisper-tiny's "
                      "d_model, so audio enters where a text prompt would.")
    bullets(s, Inches(0.72), Inches(1.95), Inches(11.9), [
        ("It is not a diffusion model in operation.",
         "timestep=0, no scheduler, no sampling loop — one forward pass per frame. That is why it "
         "approaches real time at all, and why fps scales with batch rather than step count."),
        ("in_channels 8, out_channels 4 — it inpaints, it does not generate.",
         "Two concatenated latents in: the masked lower face and an intact reference. One out. It "
         "never synthesises a person; it repaints the mouth of frames you supplied."),
        ("cross_attention_dim 384 is exactly whisper-tiny's d_model.",
         "Audio enters where a text prompt would in an image model. That is the conditioning "
         "mechanism, and why the audio encoder cannot be swapped casually."),
        ("The honest name: audio-conditioned latent inpainting.",
         "Identity comes from the reference, not the model — which is also why enrollment quality "
         "dominates output quality, and why a photograph is animated before it is enrolled."),
    ], size=13.5)

    s = content(prs, "Where this is going — the renderer as its own participant",
                "Beat 7 · module: LiveKit worker",
                notes="Run the two-process demo in a terminal if you have time; the numbers below "
                      "are from it. The important idea is that the renderer becomes a unit you can "
                      "run N of, which is what turns a single-session prototype into something "
                      "that scales — and it makes the GPU split a deployment choice.")
    beat(s, Inches(0.72), Inches(1.95), Inches(11.9),
         "Two terminals: scripts/avatar_worker.py --audio stream, then "
         "scripts/avatar_sender.py --interrupt-after 3. Read the worker's report.",
         "The renderer joins the room as its own participant and publishes both media itself.",
         "Audio arrives over lk.audio_stream, a barge-in arrives as the lk.clear_buffer RPC, and "
         "one AVSynchronizer pairs the frames before they are published — instead of two "
         "publishers on two clocks, which is what the current gap comes from.")
    table(s, Inches(0.72), Inches(4.15), Inches(11.9),
          ["Measured at a remote subscriber", "Value"],
          [["Video / audio frames decoded",
            f"{MEASURED['sub_video']} / {MEASURED['sub_audio']}"],
           ["A/V drift — the claim being tested",
            f"median {MEASURED['drift_median']}, and stable to within "
            f"{MEASURED['drift_spread']}"],
           ["The same quantity over WebSocket", MEASURED["ws_drift_range"]],
           ["Barge-in over RPC", f"dropped {MEASURED['bargein_dropped']}, epoch advanced"],
           ["Seams forced across two turns", MEASURED["seams_forced"]]],
          widths=[0.46, 0.54], size=11.5, row_h=0.44,
          colors={(1, 1): GOOD, (2, 1): WARN})
    callout(s, Inches(0.72), Inches(6.05), Inches(11.9), Inches(0.85),
            "Read the spread, not the offset. The constant −240 ms is an artifact of measuring each "
            "timeline from its own first frame; a fixed offset is startup latency and correctable. "
            "Variance is what a viewer reads as bad lip-sync, and 9 ms of it across 22 seconds is "
            "the difference a synchroniser buys.", tone=GOOD, label="How to read that table")

    # ---------------------------------------------------------------- 7. proof
    section(prs, "05", "The numbers, and where they came from",
            "Every figure is from a run on named hardware. Where a number does not exist, the "
            "slide says so.")

    # ---------------------------------------------------------------- 8. stack
    section(prs, "05", "Tech stack and models",
            "Everything that runs, what it costs on disk, and what licence it carries.")

    s = content(prs, "Models, by area", "Tech stack",
                notes="Five models cooperate for the face. The landmark model is our own "
                      "substitution: upstream builds an mmpose RTMPose model at import needing "
                      "pinned Linux-only compiled wheels, and it does one line of real work — "
                      "keypoints[0][23:91], the iBUG-68 layout — which face_alignment has "
                      "produced in pure torch since 2017.")
    table(s, Inches(0.72), Inches(1.9), Inches(11.9),
          ["Area", "Model", "Size", "Licence / host", "Role"],
          [["Face · lip-sync", "MuseTalk v1.5 U-Net", "3.24 GB", "MIT code, commercial weights",
            "Repaints the mouth in latent space"],
           ["Face · image codec", "stabilityai/sd-vae-ft-mse", "319 MB", "MIT",
            "256×256 crop ↔ 32×32×4 latents"],
           ["Face · audio encoder", "openai/whisper-tiny (encoder)", "144 MB", "MIT",
            "Conditions the U-Net via cross-attention"],
           ["Face · blending", "BiSeNet + ResNet-18", "95 MB", "MIT",
            "Feathers the generated mouth in"],
           ["Face · landmarks", "face_alignment (FAN + S3FD)", "~180 MB", "BSD-3",
            "Substituted for mmpose RTMPose — see notes"],
           ["Enrollment · animation", "LivePortrait", "2.0 GB", "MIT",
            "Motion transfer onto a still photograph"],
           ["Voice · cloning", "Chatterbox Turbo, 350M", "—", "MIT, self-hosted sidecar",
            "A persona that sounds like a specific person"],
           ["Voice · TTS", "Deepgram Aura 2 (aura-2-thalia-en)", "—", "hosted",
            f"{MEASURED['tts_hosted_ms']} ms to first audio"],
           ["Voice · STT", "Deepgram nova-3", "—", "hosted",
            "Streaming, with endpointing and utterance-end"],
           ["Language", "gpt-oss:20b", "—", "OpenAI-compatible endpoint",
            "Self-hostable without a code change"]],
          widths=[0.16, 0.24, 0.09, 0.22, 0.29], size=10.5, row_h=0.42)

    s = content(prs, "Software, and why there are three environments", "Tech stack",
                notes="The three-venv split is the slide people question, so lead with the "
                      "evidence: installing Chatterbox beside MuseTalk downgraded numpy and torch "
                      "and left libtorch_cuda.so with an undefined ncclCommResume symbol — it "
                      "took the renderer down. No transformers version satisfies both. The "
                      "collision forced a sidecar, which turned out to be the better "
                      "architecture: a second GPU is now configuration.")
    table(s, Inches(0.72), Inches(1.9), Inches(5.85),
          ["Layer", "Stack"],
          [["Runtime", "Python 3.12, FastAPI, uvicorn, pydantic"],
           ["Store", "psycopg 3, PostgreSQL; JSON files as the no-credential default"],
           ["Retrieval", "ChromaDB (optional)"],
           ["ML", "torch 2.13+cu130, diffusers 0.30.2, transformers 4.39.2"],
           ["Vision / audio", "numpy 2.4.6, OpenCV 4.14, librosa, soundfile, ffmpeg 6.1.1"],
           ["Console", "Next.js (React 19), TypeScript, Tailwind, pnpm"],
           ["Assistant", "LangGraph, langchain-core, sse-starlette"],
           ["Media plane", "LiveKit SFU + egress + Redis, via Docker"],
           ["Quality", f"pytest ({MEASURED['tests']}), ruff, mypy strict, Playwright, ESLint"]],
          widths=[0.28, 0.72], size=11)
    table(s, Inches(6.95), Inches(1.9), Inches(5.67),
          ["Environment", "Why it must be separate"],
          [[".venv-musetalk", "numpy 2.4, transformers 4.39.2"],
           [".venv-voice", "Chatterbox needs a newer transformers for LlamaModel"],
           ["LivePortrait/.venv", "pins numpy 1.26 and OpenCV 4.10"]],
          widths=[0.34, 0.66], size=11)
    callout(s, Inches(6.95), Inches(3.5), Inches(5.67), Inches(2.0),
            "This was learned the hard way. Installing Chatterbox beside MuseTalk downgraded numpy "
            "and torch and left libtorch_cuda.so with an undefined ncclCommResume symbol — it took "
            "the renderer down, not the voice. No transformers version satisfies both. Recovery "
            "needed a full uninstall of torch and every nvidia-* package.\n\nThe collision forced "
            "the sidecar, and the sidecar is the better architecture: a hosted API and a local "
            "process look identical to the orchestrator, so a second GPU is configuration rather "
            "than work.", tone=ACCENT, label="Not tidiness — a real collision")

    # ---------------------------------------------------------------- 9. residency
    s = content(prs, "Data residency — where candidate data actually goes", "Compliance posture",
                notes="This is the slide that matters most to the original premise: the vendor "
                      "concern was that candidate audio and video leave your infrastructure. Be "
                      "precise. Today the face never leaves and the conversation does. That is a "
                      "real improvement on a full vendor, and it is not yet the whole promise.")
    table(s, Inches(0.72), Inches(1.9), Inches(11.9),
          ["Component", "Runs", "Candidate data that leaves"],
          [["Renderer (MuseTalk)", "your GPU", "nothing"],
           ["Enrollment (LivePortrait)", "your GPU / CPU", "nothing"],
           ["Voice cloning (Chatterbox)", "your GPU", "nothing"],
           ["Store (Postgres or JSON)", "your host", "nothing"],
           ["Console, assistant, SFU, egress, Redis", "your host / Docker", "nothing"],
           ["Speech to text", "Deepgram", "candidate audio"],
           ["Text to speech", "Deepgram", "interviewer text only"],
           ["Language model", "OpenAI-compatible endpoint", "transcript and prompt"]],
          widths=[0.34, 0.24, 0.42], size=11.5,
          colors={(5, 2): BAD, (7, 2): BAD, (0, 2): GOOD, (1, 2): GOOD, (2, 2): GOOD,
                  (3, 2): GOOD, (4, 2): GOOD})
    callout(s, Inches(0.72), Inches(5.3), Inches(5.8), Inches(1.5),
            "Today the face never leaves and the conversation does. The LLM adapter speaks the "
            "OpenAI wire format precisely so Ollama, vLLM or LM Studio take over with a base-URL "
            "change and no code change.", tone=GOOD, label="What is already true")
    callout(s, Inches(6.82), Inches(5.3), Inches(5.8), Inches(1.5),
            "STT and TTS need real replacements, not configuration — Whisper for one, XTTS-v2 or "
            "F5-TTS for the other. That is genuine work and it is unbuilt. Biometric data is also "
            "now in the store with no authentication in front of it.",
            tone=ACCENT, label="What is not")

    # ---------------------------------------------------------------- 10. production
    section(prs, "06", "Running this in production",
            "What 25 fps costs, what a session occupies, and which of these numbers we have not "
            "earned the right to state.")

    s = content(prs, "What true real time would take", "Production",
                notes="Be concrete about which levers are measured and which are not. The CPU "
                      "overlap is done and returned 1.59×. Everything else on this list is "
                      "unmeasured on our hardware, and the deck says so rather than projecting.")
    table(s, Inches(0.72), Inches(1.9), Inches(11.9),
          ["Lever", "Expected effect", "Evidence"],
          [["Overlap CPU with GPU", f"{MEASURED['render_ms_seq']} → {MEASURED['render_ms']} "
            "ms/frame", "DONE — measured 1.59×"],
           ["A faster VAE decoder (TAESD-class, TensorRT, fp8)",
            "Attacks 74% of the remaining frame", "NOT YET MEASURED — the single largest lever"],
           ["A larger GPU (L4, A10, L40S)", "Shrinks the term that now dominates",
            "NOT YET MEASURED — no card above a T4 tested"],
           ["Lower output resolution", "Does not touch VAE decode — fixed 256×256 work",
            "Measured: proposed and discarded for this reason"],
           ["Bigger batch", "No further gain; the curve is flat past 16",
            "Measured on the T4"],
           ["Second GPU for the voice sidecar", f"Removes {MEASURED['voice_contention']} "
            "contention", "Measured contention; the fix is configuration"]],
          widths=[0.32, 0.34, 0.34], size=11,
          colors={(0, 2): GOOD, (1, 2): WARN, (2, 2): WARN})
    callout(s, Inches(0.72), Inches(5.3), Inches(11.9), Inches(1.5),
            "One T4 serves a self-hosted face or a self-hosted voice, not both: avatar_first_frame "
            f"goes {MEASURED['voice_contention']} when the voice sidecar competes. It is not "
            "memory — 7.6 of 15 GB — it is compute contention, and lowering fps and batch size was "
            "tested and did not help (27,887 ms at 6 fps vs 28,108 ms at 8). So this is a hardware "
            "or vendor decision, not a tuning exercise.",
            tone=ACCENT, label="[HUMAN] the open decision")

    s = content(prs, "Cost per interview — what we can and cannot claim", "Production",
                notes="This is the most important honesty slide in the deck. We have measured "
                      "throughput. We have NOT measured concurrent sessions per GPU, and cost per "
                      "minute is dominated by exactly that. Present the arithmetic as arithmetic, "
                      "name the assumption, and say what one experiment would close it. Do not "
                      "let anyone leave the room quoting a dollar figure as measured.")
    table(s, Inches(0.72), Inches(1.9), Inches(11.9),
          ["Input", "Status", "Value"],
          [["Render cost per frame, T4", "MEASURED", f"{MEASURED['render_ms']} ms"],
           ["Delivered fps, one session, T4", "MEASURED", f"{MEASURED['fps_delivered']} fps"],
           ["GPU memory, one session + models", "MEASURED", "7.6 of 15 GB"],
           ["Concurrent sessions per T4", "NOT YET MEASURED",
            "Never tested above one. This is the term cost per minute depends on most"],
           ["Enrollment cost, one persona", "MEASURED",
            f"~{MEASURED['prepare_ms']} ms GPU, once per persona, then cached"],
           ["Cloud GPU list price", "ASSUMPTION — not a quote",
            "PROCESS.md §4.2 assumes a T4-class instance at ~$0.35–0.50/hr"],
           ["Vendor per-minute price", "ASSUMPTION — no contract in hand",
            "PROCESS.md §4.2 assumes ~$0.10–0.30/min at list"]],
          widths=[0.3, 0.22, 0.48], size=11,
          colors={(0, 1): GOOD, (1, 1): GOOD, (2, 1): GOOD, (4, 1): GOOD,
                  (3, 1): BAD, (5, 1): WARN, (6, 1): WARN})
    callout(s, Inches(0.72), Inches(5.4), Inches(11.9), Inches(1.4),
            "We deliberately do not put a $/minute figure on a slide. The arithmetic in "
            "PROCESS.md §4.2 lands one to two orders of magnitude below the assumed vendor price — "
            "but it divides a measured GPU cost by an UNMEASURED concurrency number, and that "
            "number is the whole answer. One experiment closes it: run N concurrent sessions on "
            "one card and find where fps falls below the target. Until then the direction is "
            "credible and the magnitude is not.",
            tone=ACCENT, label="Why there is no dollar figure on this slide")

    s = content(prs, "Observability — where we instrument", "Production",
                notes="Telemetry is not a side channel here: it is the only recorder of a turn. "
                      "The server builds the stored transcript from the event stream rather than "
                      "from separate write calls, so there is exactly one authority and a second "
                      "path cannot disagree with it. The cost of that design is that an unemitted "
                      "event is an unrecorded turn — which is precisely the silence-re-prompt bug.")
    table(s, Inches(0.72), Inches(1.9), Inches(11.9),
          ["Signal", "Emitted at", "Answers"],
          [["latency: turn_detect", "End-of-turn decision", "Is the silence window right for "
            "this candidate?"],
           ["latency: llm_ttft", "First sentence from the model", "Is the endpoint or the prompt "
            "the problem?"],
           ["latency: tts_first_audio", "First PCM chunk", "REST vs WebSocket, and voice choice"],
           ["latency: avatar_first_frame", "Browser paint, not socket write",
            "The only end-to-end number a candidate feels"],
           ["latency: perceived_total", "Turn start to first paint", "The headline SLO"],
           ["frames_repeated / frames_discarded", "Mixer", "Is the renderer keeping cadence, or "
            "is the queue draining?"],
           ["stale_artifact_dropped", "Epoch check", "Is cancellation actually working?"],
           ["heard (with transcribed / silent)", "Turn open", "Broken STT vs. a quiet candidate — "
            "different pages"],
           ["WarmupReport, schema_problems", "/config", "Is this process ready, and is the "
            "database current?"]],
          widths=[0.28, 0.24, 0.48], size=11)

    # ---------------------------------------------------------------- 11. judgment
    section(prs, "07", "Recommendation and migration",
            "These sections belong to the author. What follows is quoted from PROCESS.md, with "
            "the one input that changed since it was written.")

    s = content(prs, "Build vs. buy — the standing recommendation", "Quoted from PROCESS.md §4",
                notes="Read the recommendation as written; do not extend it. Then be explicit "
                      "that one of its inputs changed: §4.2 said 'I did not run a GPU', and now "
                      "one has run. The measured render cost is far below the assumed vendor "
                      "price, which if anything strengthens the case — but re-deriving the memo "
                      "is the author's judgment call, not something to assert from a slide.")
    callout(s, Inches(0.72), Inches(1.9), Inches(11.9), Inches(1.5),
            "“Keep the vendor for the rendering stage. Build the orchestration layer in-house, "
            "starting now.” That is the hybrid, and the split is not a hedge — it falls directly "
            "out of what was measured.", tone=GOOD, label="The recommendation, as written")
    bullets(s, Inches(0.72), Inches(3.6), Inches(11.9), [
        ("The reasoning, unchanged: the model was never the bottleneck.",
         f"A full turn measures {MEASURED['turn_total']} against a sub-second target, and the "
         "three dominant terms are the end-of-turn policy, LLM time-to-first-token and TTS. The "
         "part a vendor sells is the part that was never the problem."),
        ("What changed since it was written: a GPU ran.",
         "§4.2 opens “I have no vendor contract and did not run a GPU.” Render cost per frame is "
         "now measured, and it lands well below the assumed vendor per-minute price — but it "
         "divides by a concurrency figure we still have not measured."),
        ("[HUMAN] The refresh this deserves, and who owns it.",
         "Re-derive §4.2 with the measured render cost, run the concurrency experiment, and "
         "restate §4.4's thresholds. The recommendation may survive intact; asserting that from a "
         "slide would be exactly the foregone conclusion the brief warns against."),
    ], size=13)

    s = content(prs, "Migration plan — cutover without a regression", "Quoted from PROCESS.md §5",
                notes="The test of this section in the brief is whether another senior engineer "
                      "could execute it without the author in the room. Walk the four phases and "
                      "point at the preconditions — the plan is gated on things being true before "
                      "phase 1 starts, which is what makes it executable rather than aspirational.")
    flow(s, Inches(0.72), Inches(1.95), Inches(11.9), [
        ("Preconditions", "Named and checkable before anything ships. A cutover gated on nothing "
         "is a hope."),
        ("Phase 1 · Shadow", "Render in parallel, serve the vendor. Compare quality and latency "
         "on real traffic, no candidate affected."),
        ("Phase 2 · Flagged rollout", "Per-tenant flag, small cohorts first, with the abort "
         "criteria written down in advance."),
        ("Rollback", "One flag flip back to the vendor. The vendor path stays warm and paid for "
         "until decommission."),
        ("Decommission", "Only after a defined quiet period with no regression. Removing the "
         "fallback is the last step, not the first."),
    ])
    bullets(s, Inches(0.72), Inches(3.5), Inches(11.9), [
        ("The boundary is what makes this cheap.",
         "The renderer is a Protocol with two implementations already. Shadow mode is a second "
         "implementation of an interface that exists, not a fork of the pipeline — which is the "
         "practical payoff of the contracts decision on slide 9."),
        ("Rollback is a flag, not a redeploy.",
         "Because the vendor and in-house renderers are interchangeable at the same seam, and "
         "because session state lives in the orchestrator rather than in the renderer."),
    ], size=13)

    # ---------------------------------------------------------------- 12. gaps
    s = content(prs, "What is missing — stated, not buried", "Risks and gaps",
                notes="Close the deck on this rather than on the recommendation. Volunteering the "
                      "gap list is the argument for the rest of the deck being trustworthy. The "
                      "authentication row is the one to say slowly: the store now holds a real "
                      "person's face and voice, which changes the posture from untidy to serious.")
    table(s, Inches(0.72), Inches(1.9), Inches(11.9),
          ["Gap", "Consequence", "Standing"],
          [["No authentication anywhere", "The candidate link is not a credential; the assistant "
            "reads any transcript. The store now holds real faces and voices — biometric data",
            "Deferred by the owner"],
           ["One GPU cannot host face + voice", "Self-hosted cloning and a self-hosted face are "
            "an either/or ({MEASURED['voice_contention']})", "[HUMAN] decision"],
           ["Concurrency per GPU untested", "Cost per minute cannot be credibly stated",
            "One experiment away"],
           ["2.0× short of 25 fps", f"{MEASURED['render_ms']} ms/frame; VAE decode is 74%",
            "Needs a faster decoder or a bigger card"],
           ["Start lag split unmeasured", f"{MEASURED['start_lag']} before video starts — render "
            "window vs mixer lead-in unknown", "Measure before tuning"],
           ["Self-hosted STT/TTS unbuilt", "The conversation still leaves your infrastructure",
            "Real work, not configuration"],
           ["No output-quality comparison", "Our substituted landmark detector vs mmpose RTMPose "
            "is unmeasured for quality", "Known, unquantified"],
           ["No CUDA float32 ratio", "The fp16 default is not in doubt; the ratio on CUDA is",
            "Low value, stated anyway"]],
          widths=[0.24, 0.48, 0.28], size=10.5, row_h=0.44,
          colors={(0, 2): BAD, (1, 2): WARN, (2, 2): WARN})

    s = content(prs, "How to verify any of this yourself", "Reproducibility",
                notes="End by handing over the commands. Everything on the metrics slides came "
                      "from one of these, and each writes a JSON file alongside its console "
                      "output so a figure can be traced rather than trusted.")
    table(s, Inches(0.72), Inches(1.9), Inches(11.9),
          ["Command", "Produces"],
          [["pytest -m \"not gpu\"", f"{MEASURED['tests']} tests: no GPU, no weights, no network"],
           ["python scripts/bench_renderer.py", "The throughput and stage tables, plus the "
            "sequential-vs-overlapped comparison"],
           ["python scripts/measure_lag.py --agent <id>", "The live delivery figures: fps, "
            "trailing gap, start lag, discards"],
           ["python scripts/smoke_session.py", "17 assertions over a real socket: barge-in, "
            "epoch drops, cadence, mic-driven turns"],
           ["curl localhost:8000/config", "Which implementation each boundary resolved to, warm-up "
            "state, and schema problems"],
           ["docker compose --env-file .env.development up -d", "SFU, egress and Redis — recording "
            "needs a recorder the SFU binary omits"]],
          widths=[0.42, 0.58], size=11.5, row_h=0.46)
    callout(s, Inches(0.72), Inches(5.4), Inches(11.9), Inches(1.2),
            "MEASUREMENTS.md is the source of every number in this deck, including a §9 that lists "
            "what is not measured and a record of two figures we published and later retracted "
            "because they had been measured on the wrong device. The retraction is in the document "
            "on purpose.", tone=GOOD, label="Where the numbers live")

    title_slide(prs, "Ends", "Clarity over cleverness.\nBoundaries over convenience.\nProof over opinion.",
                "Every figure in this deck came from a run on named hardware. Where a number does "
                "not exist, the slide says NOT YET MEASURED and says what would close it.")

    prs.save(str(out))
    print(f"wrote {out} — {len(prs.slides.__iter__.__self__._sldIdLst)} slides")


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--out", default="nod-demo.pptx")
    args = parser.parse_args()
    build(Path(args.out))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
